from pptx import Presentation
import sys

def extract_text_from_ppt(filepath, outpath):
    try:
        prs = Presentation(filepath)
        with open(outpath, 'w', encoding='utf-8') as f:
            for i, slide in enumerate(prs.slides):
                f.write(f"--- Slide {i+1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        f.write(shape.text + "\n")
    except Exception as e:
        print(f"Error reading {filepath}: {e}")

if __name__ == "__main__":
    filepath = sys.argv[1]
    outpath = sys.argv[2]
    extract_text_from_ppt(filepath, outpath)
